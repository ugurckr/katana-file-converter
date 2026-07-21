"""Doküman dönüştürücüleri: pdf, docx, pptx, epub, mobi. Tümü saf Python;
docx/pptx->pdf ve mobi->epub artık harici araç (LibreOffice/Calibre) gerektirmez."""

import io
import os
import shutil
import tempfile
from pathlib import Path

import ebooklib
import fitz
import mammoth
import reportlab
from ebooklib import epub
from lxml import html
from mobi.kindleunpack import unpackBook
from PIL import Image
from pdf2docx import Converter
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas as rl_canvas
from xhtml2pdf import pisa
from xhtml2pdf.default import DEFAULT_FONT

from .base import register

# 914400 EMU/inç ÷ 72 punto/inç — python-pptx EMU'larını PDF puntosuna çevirir.
_EMU_PER_PT = 12700

# Standart PDF fontu (Helvetica) Türkçe'ye özgü ş/ğ/ı/İ glifleri içermez. reportlab'ın
# paketlediği Bitstream Vera Sans tam Unicode kapsar; hem reportlab (pptx) hem xhtml2pdf
# (docx/epub) çıktıları için kaydedilir.
_FONTS_DIR = os.path.join(os.path.dirname(reportlab.__file__), "fonts")
_FONT_NAME = "KatanaSans"
_FONT_BOLD = "KatanaSans-Bold"
_FONT_CSS = "katanasans"  # xhtml2pdf CSS'inde kullanılan font-family adı (küçük harf)
_FONT_PATH = os.path.join(_FONTS_DIR, "Vera.ttf")
pdfmetrics.registerFont(TTFont(_FONT_NAME, _FONT_PATH))
pdfmetrics.registerFont(TTFont(_FONT_BOLD, os.path.join(_FONTS_DIR, "VeraBd.ttf")))
# xhtml2pdf'in font-family eşlemesine ekle; @font-face temp-kopyalama yolunu (bu ortamda
# Türkçe yol karakterleriyle bozuluyor) atlayarak reportlab'a kayıtlı fontu doğrudan kullanır.
DEFAULT_FONT[_FONT_CSS] = _FONT_NAME

# pptx renkleri açık/koyu zemine göre okunur kalsın diye varsayılan metin renkleri (0-1 RGB).
_TEXT_ON_LIGHT = (0.10, 0.10, 0.10)
_TEXT_ON_DARK = (0.92, 0.92, 0.92)


def _pdf_to_images(src: Path, dst: Path, fmt: str) -> None:
    dst.parent.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(src)
    try:
        if doc.page_count == 0:
            raise ValueError(f"'{src}' içinde sayfa bulunamadı.")
        for i, page in enumerate(doc):
            pix = page.get_pixmap()
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            if fmt == "JPEG":
                img = img.convert("RGB")
            out_path = dst if i == 0 else dst.with_name(f"{dst.stem}__p{i + 1:02d}{dst.suffix}")
            img.save(out_path, format=fmt)
    finally:
        doc.close()


def _html_to_pdf(html_body: str, src: Path, dst: Path) -> None:
    """xhtml2pdf ile HTML gövdesini PDF'e basar (epub_to_pdf ile aynı desen).
    Türkçe karakterlerin doğru görünmesi için gövde fontu Vera'ya (KatanaSans) ayarlanır."""
    style = f'<style> html, body {{ font-family: {_FONT_CSS}; }}</style>'
    document = '<!doctype html>\n<meta charset="utf-8">\n' + style + html_body
    dst.parent.mkdir(parents=True, exist_ok=True)
    with dst.open("wb") as f:
        result = pisa.CreatePDF(document, dest=f)
    if result.err:
        raise RuntimeError(f"'{src}' PDF'e dönüştürülürken hata oluştu.")


def _emu_to_pt(value: int) -> float:
    return value / _EMU_PER_PT


def _is_dark(rgb: tuple[int, int, int]) -> bool:
    """(r,g,b 0-255) rengin göz için koyu olup olmadığını luminance ile belirler."""
    r, g, b = rgb
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) < 128


@register(".pdf", ".docx", "Düzenlenebilir Word belgesi")
def pdf_to_docx(src: Path, dst: Path) -> None:
    dst.parent.mkdir(parents=True, exist_ok=True)
    cv = Converter(str(src))
    try:
        cv.convert(str(dst))
    finally:
        cv.close()


@register(".pdf", ".png", "PNG görsel (sayfa başına)", multi_output=True)
def pdf_to_png(src: Path, dst: Path) -> None:
    _pdf_to_images(src, dst, fmt="PNG")


@register(".pdf", ".jpg", "JPEG görsel (sayfa başına)", multi_output=True)
def pdf_to_jpg(src: Path, dst: Path) -> None:
    _pdf_to_images(src, dst, fmt="JPEG")


@register(".docx", ".pdf", "PDF belge")
def docx_to_pdf(src: Path, dst: Path) -> None:
    with src.open("rb") as f:
        body = mammoth.convert_to_html(f).value
    _html_to_pdf(body, src, dst)


def _solid_rgb(fill) -> tuple[int, int, int] | None:
    """Bir FillFormat düz (solid) RGB ise (r,g,b) döner; değilse/tema rengi ise None."""
    try:
        if fill.type != MSO_FILL.SOLID:
            return None
        rgb = fill.fore_color.rgb  # tema rengiyse burada hata verir
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def _slide_bg_rgb(slide) -> tuple[int, int, int] | None:
    """Slaytın efektif arka plan rengini slide → layout → master zincirinde arar."""
    for source in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        rgb = _solid_rgb(source.background.fill)
        if rgb is not None:
            return rgb
    return None


def _run_size_pt(run) -> float:
    size = run.font.size
    return size.pt if size is not None else 18.0


def _run_color(run, default: tuple[float, float, float]) -> tuple[float, float, float]:
    """Run'ın açık RGB rengini (0-1) döner; tema rengi/renksizse varsayılana düşer."""
    try:
        color = run.font.color
        if color.type is not None:
            rgb = color.rgb  # tema rengiyse hata verir
            return (rgb[0] / 255, rgb[1] / 255, rgb[2] / 255)
    except Exception:
        pass
    return default


def _draw_text_frame(pdf, shape, page_h: float, default_color: tuple[float, float, float]) -> None:
    """Bir metin çerçevesini run bazında renk/kalınlık ve paragraf hizasıyla çizer."""
    x0 = _emu_to_pt(shape.left)
    box_w = _emu_to_pt(shape.width)
    cursor = page_h - _emu_to_pt(shape.top)

    for para in shape.text_frame.paragraphs:
        runs = [r for r in para.runs if r.text]
        line_h = (max((_run_size_pt(r) for r in runs), default=18.0)) * 1.2
        cursor -= line_h
        if not runs:
            continue

        segments = []
        total_w = 0.0
        for run in runs:
            size = _run_size_pt(run)
            font = _FONT_BOLD if run.font.bold else _FONT_NAME
            width = pdfmetrics.stringWidth(run.text, font, size)
            segments.append((run.text, font, size, _run_color(run, default_color), width))
            total_w += width

        if para.alignment == PP_ALIGN.CENTER:
            x = x0 + (box_w - total_w) / 2
        elif para.alignment == PP_ALIGN.RIGHT:
            x = x0 + box_w - total_w
        else:
            x = x0

        for text, font, size, color, width in segments:
            pdf.setFont(font, size)
            pdf.setFillColorRGB(*color)
            pdf.drawString(x, cursor, text)
            x += width


@register(".pptx", ".pdf", "PDF belge")
def pptx_to_pdf(src: Path, dst: Path) -> None:
    """Slaytları reportlab ile PDF'e çizer: arka plan rengi, run bazlı metin rengi/kalınlığı
    ve paragraf hizası korunur. Tablolar, gruplanmış şekiller ve animasyonlar çizilmez."""
    prs = Presentation(str(src))
    page_w = _emu_to_pt(prs.slide_width)
    page_h = _emu_to_pt(prs.slide_height)

    dst.parent.mkdir(parents=True, exist_ok=True)
    pdf = rl_canvas.Canvas(str(dst), pagesize=(page_w, page_h))
    for slide in prs.slides:
        bg = _slide_bg_rgb(slide)
        if bg is not None:
            pdf.setFillColorRGB(bg[0] / 255, bg[1] / 255, bg[2] / 255)
            pdf.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        # Zemin koyuysa metin varsayılanı açık olsun (aksi halde okunmaz).
        default_color = _TEXT_ON_DARK if (bg is not None and _is_dark(bg)) else _TEXT_ON_LIGHT

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = ImageReader(io.BytesIO(shape.image.blob))
                    w = _emu_to_pt(shape.width)
                    h = _emu_to_pt(shape.height)
                    x = _emu_to_pt(shape.left)
                    # reportlab sol-alt orijinli; pptx sol-üst orijinli — y çevrilir.
                    y = page_h - _emu_to_pt(shape.top) - h
                    pdf.drawImage(image, x, y, width=w, height=h, mask="auto")
                except Exception:
                    continue
            elif shape.has_text_frame:
                _draw_text_frame(pdf, shape, page_h, default_color)
        pdf.showPage()
    pdf.save()


@register(".epub", ".pdf", "PDF belge")
def epub_to_pdf(src: Path, dst: Path) -> None:
    book = epub.read_epub(str(src))
    chapters = [
        item.get_content().decode("utf-8", errors="ignore")
        for item in book.get_items()
        if item.get_type() == ebooklib.ITEM_DOCUMENT
    ]
    _html_to_pdf("<hr/>\n".join(chapters), src, dst)


@register(".pdf", ".txt", "Düz metin (metin çıkarma)")
def pdf_to_txt(src: Path, dst: Path) -> None:
    doc = fitz.open(src)
    try:
        text = "\n\n".join(page.get_text().strip() for page in doc)
    finally:
        doc.close()
    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(text.strip() + "\n", encoding="utf-8")


@register(".epub", ".txt", "Düz metin (metin çıkarma)")
def epub_to_txt(src: Path, dst: Path) -> None:
    book = epub.read_epub(str(src))
    chapters = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            tree = html.fromstring(item.get_content())
            chapters.append(tree.text_content().strip())
    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text("\n\n".join(filter(None, chapters)) + "\n", encoding="utf-8")


@register(".mobi", ".epub", "EPUB e-kitap")
def mobi_to_epub(src: Path, dst: Path) -> None:
    """mobi paketi (KindleUnpack fork'u) mobi'yi açar. Modern KF8/AZW3 içerik
    mobi8/*.epub verir; eski MOBI6 epub üretemez ve anlaşılır hata döndürür.
    mobi.extract() hata halinde kendi temp dizinini sızdırdığı için, temizliği
    garanti etmek adına unpackBook doğrudan kendi TemporaryDirectory'mizle çağrılır."""
    with tempfile.TemporaryDirectory(prefix="katana-mobi") as tempdir:
        try:
            unpackBook(str(src), tempdir, epubver="A")
        except Exception as exc:
            # unpackBook bozuk/geçersiz dosyalarda teknik (struct) hatalar fırlatır;
            # kullanıcıya anlaşılır bir mesaja çevir.
            raise RuntimeError(
                f"'{src}' geçerli bir MOBI dosyası değil ya da okunamadı."
            ) from exc

        base = src.stem
        produced = Path(tempdir) / "mobi8" / f"{base}.epub"
        if not produced.exists():
            raise RuntimeError(
                f"'{src}' eski bir MOBI (MOBI6) formatında; EPUB'a dönüştürülemedi."
            )
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(produced, dst)